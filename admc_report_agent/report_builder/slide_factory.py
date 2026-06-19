from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
import os
import glob as _glob

# Base directory for resolving input paths relative to the package
_PKG_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# Support both "input" and "Inputs" folder names — check both because
# git may create an empty "input/" while user files live in "Inputs/"
_INPUT_DIRS = []
for _candidate in ("input", "Inputs"):
    _d = os.path.join(_PKG_DIR, _candidate)
    if os.path.isdir(_d):
        _INPUT_DIRS.append(_d)
if not _INPUT_DIRS:
    _INPUT_DIRS = [os.path.join(_PKG_DIR, "input")]

COLORS = {
    "teal_primary": "0097B2",
    "teal_dark": "007A91",
    "teal_light": "E8F7FA",
    "teal_mid": "B3E0EA",
    "white": "FFFFFF",
    "navy_text": "1A2B4A",
    "dark_title": "0D1B2E",
    "gray_light": "F5F5F5",
    "gray_border": "CCCCCC",
}

FONT = "Montserrat"
SLIDE_W = 13.33
SLIDE_H = 7.5


def hex_to_rgb(hex_color):
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def add_text_box(slide, text, left, top, width, height,
                 font_name=FONT, font_size=14,
                 bold=False, color="FFFFFF", align=PP_ALIGN.LEFT,
                 word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=14, color="1A2B4A", bullet_color="0097B2"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.level = 0
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = hex_to_rgb(color)
    return txBox


def add_footer(slide, client_name):
    client_name = client_name or "client"
    add_text_box(slide, "www.activedmc.com",
                 0.3, SLIDE_H - 0.4, 3, 0.3,
                 font_size=9, color="999999")
    email = f"{client_name.lower().replace(' ', '')}@activedmc.com"
    add_text_box(slide, email,
                 SLIDE_W - 3.3, SLIDE_H - 0.4, 3, 0.3,
                 font_size=9, color="999999", align=PP_ALIGN.RIGHT)


def add_header_rule(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.7), Inches(SLIDE_W), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb("0097B2")
    shape.line.fill.background()


def add_teal_background(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb("0097B2")


def _add_logo_scaled(slide, img_path, left, top, max_w=2.0, max_h=0.7):
    """Add a logo to the slide, scaled to fit within max_w x max_h while preserving aspect ratio."""
    from PIL import Image
    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
        ratio = img_w / img_h
        w = max_w
        h = w / ratio
        if h > max_h:
            h = max_h
            w = h * ratio
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(w), Inches(h))
    except Exception:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(max_w), Inches(max_h))


def _find_clients_dirs():
    """Return all existing client logo directories across input/Inputs/assets."""
    dirs = []
    for base in _INPUT_DIRS:
        for sub in ("logos/clients", "logos/Clients"):
            d = os.path.join(base, sub)
            if os.path.isdir(d):
                dirs.append(d)
    assets_d = os.path.join(_PKG_DIR, "assets", "clients")
    if os.path.isdir(assets_d):
        dirs.append(assets_d)
    return dirs


def load_logo(client_name, assets_dir=None):
    """Find a client logo image, searching all known directories."""
    if not client_name:
        return None
    extensions = [".png", ".jpg", ".jpeg", ".webp"]
    variants = [
        client_name,
        client_name.lower(),
        client_name.upper(),
        client_name.replace(" ", "_"),
        client_name.replace(" ", "-"),
        client_name.lower().replace(" ", "_"),
        client_name.lower().replace(" ", "-"),
    ]
    client_lower = client_name.lower().split()[0] if client_name.strip() else ""

    for clients_dir in _find_clients_dirs():
        for name in variants:
            for ext in extensions:
                path = os.path.join(clients_dir, name + ext)
                if os.path.exists(path):
                    return path
        if client_lower:
            for fname in os.listdir(clients_dir):
                if fname.startswith("."):
                    continue
                stem = os.path.splitext(fname)[0].lower()
                ext_lower = os.path.splitext(fname)[1].lower()
                if ext_lower in extensions and (client_lower in stem or stem in client_name.lower()):
                    return os.path.join(clients_dir, fname)

    if assets_dir:
        for name in variants:
            for ext in extensions:
                path = os.path.join(assets_dir, "clients", name + ext)
                if os.path.exists(path):
                    return path

    return None


def _find_admc_dirs():
    """Return all existing ADMC logo directories."""
    dirs = []
    for base in _INPUT_DIRS:
        for sub in ("logos/active_dmc", "logos/Active_DMC", "logos/active_DMC",
                     "logos/ActiveDMC", "logos/ACTIVE_DMC"):
            d = os.path.join(base, sub)
            if os.path.isdir(d):
                dirs.append(d)
    return dirs


def load_admc_logo():
    """Find the Active DMC logo image, searching all known directories."""
    image_exts = ("*.png", "*.jpg", "*.jpeg", "*.webp")
    for admc_dir in _find_admc_dirs():
        for candidate in ["active_dmc_logo.png", "logo.png"]:
            path = os.path.join(admc_dir, candidate)
            if os.path.exists(path):
                return path
        for ext_pattern in image_exts:
            matches = _glob.glob(os.path.join(admc_dir, ext_pattern))
            if matches:
                return matches[0]
    return None


def load_cover_image():
    """Find a cover image in input/images/ or Inputs/images/."""
    for base in _INPUT_DIRS:
        images_dir = os.path.join(base, "images")
        for name in ["cover.png", "cover.jpg"]:
            path = os.path.join(images_dir, name)
            if os.path.exists(path):
                return path
    return None


def build_cover_slide(prs, blank, client, period):
    slide = prs.slides.add_slide(blank)
    add_teal_background(slide)
    add_text_box(slide, client.upper(), 0.8, 1.5, 6, 1.2,
                 font_size=44, bold=True, color="FFFFFF")
    add_text_box(slide, "COMMUNICATIONS REPORT", 0.8, 2.8, 6, 0.6,
                 font_size=18, color="FFFFFF")
    add_text_box(slide, period, 0.8, 3.5, 6, 0.5,
                 font_size=16, color="FFFFFF")
    add_text_box(slide, "By Active Digital Marketing Communications", 0.8, 4.3, 6, 0.5,
                 font_size=12, color="FFFFFF")

    # Cover image or placeholder rectangle
    cover_img = load_cover_image()
    if cover_img:
        slide.shapes.add_picture(
            cover_img,
            Inches(7.5), Inches(0.8), Inches(5.2), Inches(5.5)
        )
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.5), Inches(0.8), Inches(5.2), Inches(5.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb("007A91")
        shape.line.fill.background()
        add_text_box(slide, "[Professional Photo]", 8.5, 3.2, 3, 0.5,
                     font_size=12, color="B3E0EA", align=PP_ALIGN.CENTER)

    # Active DMC logo or text fallback
    admc_logo = load_admc_logo()
    if admc_logo:
        _add_logo_scaled(slide, admc_logo, 0.8, 6.2, max_w=2.0, max_h=0.7)
    else:
        add_text_box(slide, "Active DMC", 0.8, 6.5, 2, 0.4,
                     font_size=11, bold=True, color="FFFFFF")

    # Client logo or text fallback
    client_logo = load_logo(client)
    if client_logo:
        _add_logo_scaled(slide, client_logo, SLIDE_W - 2.8, 6.2, max_w=2.0, max_h=0.7)
    else:
        add_text_box(slide, client, SLIDE_W - 3, 6.5, 2.5, 0.4,
                     font_size=11, bold=True, color="FFFFFF", align=PP_ALIGN.RIGHT)


def build_contents_slide(prs, blank, client, sections):
    slide = prs.slides.add_slide(blank)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(5), Inches(SLIDE_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb("0097B2")
    shape.line.fill.background()
    add_text_box(slide, "CONTENTS", 1, 2.5, 3, 1,
                 font_size=36, bold=True, color="FFFFFF")
    content_sections = [
        "Campaign Overview", "Coverage Highlights", "Deliverables Overview",
        "Coverage Analysis", "Top Publications", "Media Relations",
        "Driving the Narrative", "Observation / Feedback",
        "Strategic Recommendations", "90-Day Plan"
    ]
    y = 1.0
    for i, s in enumerate(content_sections, 1):
        add_text_box(slide, f"›  {i}. {s}", 5.5, y, 7, 0.4,
                     font_size=14, color="1A2B4A")
        y += 0.55
    add_footer(slide, client)


def build_section_divider(prs, blank, title):
    slide = prs.slides.add_slide(blank)
    add_teal_background(slide)
    for y_off in [1.0, 1.35, 1.7]:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y_off), Inches(0.15), Inches(0.15)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb("FFFFFF")
        dot.line.fill.background()
    add_text_box(slide, title.upper(), 0.8, 2.8, 11, 1.5,
                 font_size=36, bold=True, color="FFFFFF")


def build_kpi_slide(prs, blank, client, tracker_data, insights):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Campaign Overview", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    metrics = [
        (str(tracker_data.get("total_placements", 0)), "Total Placements"),
        (str(tracker_data.get("unique_pr_campaigns", 0)), "PR Campaigns"),
        (str(tracker_data.get("tier1_count", 0)), "Tier 1 Placements"),
        (f"{int(tracker_data.get('total_reach', 0) or 0):,}", "Total Reach"),
    ]
    for i, (value, label) in enumerate(metrics):
        left = 0.5 + i * 3.1
        top = 1.2
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(2.8), Inches(1.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb("FFFFFF")
        card.line.color.rgb = hex_to_rgb("0097B2")
        card.line.width = Pt(1.5)
        add_text_box(slide, value, left + 0.1, top + 0.15, 2.6, 0.9,
                     font_size=48, bold=True, color="0097B2", align=PP_ALIGN.CENTER)
        add_text_box(slide, label, left + 0.1, top + 1.0, 2.6, 0.45,
                     font_size=13, color="888888", align=PP_ALIGN.CENTER)
    insight_items = []
    if isinstance(insights, dict):
        for ins in insights.get("insights", []):
            if isinstance(ins, dict):
                what = ins.get('what', '')
                implication = ins.get('implication', '')
                if len(implication) > 180:
                    implication = implication[:177] + "..."
                insight_items.append(f"{what} — {implication}")
            elif isinstance(ins, str):
                insight_items.append(ins)
    if insight_items:
        add_bullet_list(slide, insight_items[:4], 0.5, 3.2, 12, 3.5, font_size=11)
    add_footer(slide, client)


def build_highlights_slide(prs, blank, client, tracker_data):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Coverage Highlights", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    rows = tracker_data.get("rows", [])
    y = 1.2
    for r in rows[:8]:
        headline = str(r.get("headline", ""))
        pub = str(r.get("publication", ""))
        tier = str(r.get("tier", ""))
        lang = str(r.get("language", ""))
        tier_str = f"Tier {tier}" if tier and not tier.lower().startswith("tier") else tier
        line = f"{pub} — {headline} [{tier_str}, {lang}]"
        add_text_box(slide, f"• {line}", 0.5, y, 12, 0.35,
                     font_size=11, color="1A2B4A")
        y += 0.45
    add_footer(slide, client)


def build_deliverables_slide(prs, blank, client, deliverables, note=None):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Deliverables Overview", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    if not deliverables:
        msg = note or "[No deliverables data available]"
        add_text_box(slide, msg, 0.5, 2, 12, 2,
                     font_size=14, color="888888")
        add_footer(slide, client)
        return
    deliverables = deliverables[:10]
    rows_count = len(deliverables) + 1
    cols = 4
    col_widths = [1.8, 1.0, 1.0, 9.0]
    table_shape = slide.shapes.add_table(
        rows_count, cols,
        Inches(0.5), Inches(1.2),
        Inches(sum(col_widths)), Inches(rows_count * 0.55)
    )
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    headers = ["Activity", "Agreed", "Delivered", "Title"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb("0097B2")
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = hex_to_rgb("FFFFFF")
        run.font.name = FONT
    for i, item in enumerate(deliverables):
        row_idx = i + 1
        row_bg = "E8F7FA" if (i % 2 == 1) else "FFFFFF"
        titles_text = item.get("name", "")
        if item.get("items"):
            titles_text = "; ".join(t.get("name", "") for t in item["items"][:3])
        row_data = [
            item.get("type", "Other").replace("_", " ").title(),
            f"x{item.get('planned', 0)}",
            f"x{item.get('delivered', 0)}",
            titles_text[:100],
        ]
        for j, val in enumerate(row_data):
            cell = table.cell(row_idx, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(row_bg)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(11)
            run.font.color.rgb = hex_to_rgb("1A2B4A")
            run.font.name = FONT
    add_footer(slide, client)


def build_analysis_slide(prs, blank, client, tracker_data, insights):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Coverage Analysis", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    tier_data = tracker_data.get("tier_breakdown", {})
    lang_data = tracker_data.get("language_breakdown", {})
    media_data = tracker_data.get("media_type_breakdown", {})
    charts = [
        ("BY TIER", tier_data, 0.3),
        ("BY LANGUAGE", lang_data, 4.5),
        ("BY MEDIA TYPE", media_data, 8.7),
    ]
    teal_palette = ["0097B2", "66C5D6", "B3E0EA", "004F61", "CCE9F0"]
    for title, data_dict, left in charts:
        if not data_dict:
            add_text_box(slide, f"{title}\n[No data]", left, 1.2, 3.8, 2.5,
                         font_size=12, color="888888", align=PP_ALIGN.CENTER)
            continue
        add_text_box(slide, title, left, 0.9, 3.8, 0.35,
                     font_size=12, bold=True, color="0097B2", align=PP_ALIGN.CENTER)
        chart_data = ChartData()
        chart_data.categories = list(data_dict.keys())
        chart_data.add_series("Coverage", list(data_dict.values()))
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(left), Inches(1.3), Inches(3.8), Inches(2.8),
            chart_data
        )
        chart = chart_frame.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_value = True
        data_labels.show_category_name = False
        data_labels.show_percentage = False
        data_labels.font.size = Pt(9)
        data_labels.font.bold = True
        data_labels.font.color.rgb = hex_to_rgb("0D1B2E")
        series = chart.series[0]
        for i, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = hex_to_rgb(teal_palette[i % len(teal_palette)])
    insight_items = []
    if isinstance(insights, dict):
        for ins in insights.get("insights", []):
            if isinstance(ins, dict):
                what = ins.get('what', '')
                implication = ins.get('implication', '')
                if len(implication) > 200:
                    implication = implication[:197] + "..."
                insight_items.append(f"{what} — {implication}")
    if insight_items:
        add_bullet_list(slide, insight_items[:3], 0.5, 4.5, 12, 2.5, font_size=10)
    add_footer(slide, client)


def build_top_publications_slide(prs, blank, client, tracker_data):
    """Build a horizontal bar chart of the top publications by placement count."""
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Top Publications", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    top_pubs = tracker_data.get("top_publications", [])
    if not top_pubs:
        add_text_box(slide, "No publication data available.", 0.5, 2, 12, 1,
                     font_size=14, color="888888")
        add_footer(slide, client)
        return
    chart_data = ChartData()
    pub_names = [p[0] for p in top_pubs]
    pub_counts = [p[1] for p in top_pubs]
    chart_data.categories = pub_names
    chart_data.add_series("Placements", pub_counts)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5), Inches(1.2), Inches(12), Inches(4.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = False
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.category_axis.tick_labels.font.name = FONT
    chart.category_axis.tick_labels.font.color.rgb = hex_to_rgb("1A2B4A")
    chart.category_axis.format.line.fill.background()
    plot = chart.plots[0]
    plot.gap_width = 80
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.font.size = Pt(11)
    data_labels.font.bold = True
    data_labels.font.color.rgb = hex_to_rgb("FFFFFF")
    data_labels.number_format = '0'
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = hex_to_rgb("0097B2")
    total = tracker_data.get("total_placements", 0)
    if total:
        add_text_box(slide, f"Top {len(top_pubs)} of {total} total placements",
                     0.5, 5.9, 6, 0.3, font_size=10, color="888888")
    add_footer(slide, client)


def build_media_relations_slide(prs, blank, client, insights):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Media Relations & Engagement", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(5), Inches(5.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb("E8F7FA")
    shape.line.fill.background()
    add_text_box(slide, "[Media Relations Photo]", 1.5, 3.5, 3, 0.5,
                 font_size=12, color="0097B2", align=PP_ALIGN.CENTER)
    insight_items = []
    if isinstance(insights, dict):
        for ins in insights.get("insights", []):
            if isinstance(ins, dict):
                what = ins.get("what", "")
                implication = ins.get("implication", "")
                if len(implication) > 200:
                    implication = implication[:197] + "..."
                insight_items.append(f"{what} — {implication}")
    if insight_items:
        add_bullet_list(slide, insight_items[:5], 6, 1.5, 6.5, 5, font_size=11)
    else:
        add_text_box(slide, "[Media relations insights will appear here]", 6, 2, 6, 1,
                     font_size=12, color="888888")
    add_footer(slide, client)


def build_narrative_slide(prs, blank, client, insights):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Driving the Narrative", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    insight_items = []
    if isinstance(insights, dict):
        summary = insights.get("summary", "")
        if summary:
            if len(summary) > 500:
                summary = summary[:497] + "..."
            add_text_box(slide, summary, 0.5, 1.0, 12, 0.8,
                         font_size=12, color="1A2B4A")
        for ins in insights.get("insights", []):
            if isinstance(ins, dict):
                what = ins.get('what', '')
                implication = ins.get('implication', '')
                if len(implication) > 200:
                    implication = implication[:197] + "..."
                insight_items.append(f"{what} → {implication}")
    if insight_items:
        add_bullet_list(slide, insight_items[:4], 0.5, 2.0, 12, 4.8, font_size=11)
    add_footer(slide, client)


def build_observation_slide(prs, blank, client, insights):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Observation & Feedback", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.0), Inches(5.5), Inches(5.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb("E8F7FA")
    shape.line.fill.background()
    add_text_box(slide, "What Went Well", 6.5, 1.0, 6, 0.4,
                 font_size=16, bold=True, color="0097B2")
    went_well = []
    improve = []
    if isinstance(insights, dict):
        went_well = [s[:250] + "..." if len(str(s)) > 250 else str(s) for s in insights.get("went_well", [])]
        improve = [s[:250] + "..." if len(str(s)) > 250 else str(s) for s in insights.get("improve", [])]
    if went_well:
        add_bullet_list(slide, went_well[:3], 6.5, 1.5, 6, 2, font_size=10, bullet_color="0097B2")
    add_text_box(slide, "What Could Be Better", 6.5, 3.8, 6, 0.4,
                 font_size=16, bold=True, color="0097B2")
    if improve:
        add_bullet_list(slide, improve[:3], 6.5, 4.3, 6, 2.5, font_size=10, bullet_color="0097B2")
    add_footer(slide, client)


def build_recommendations_slide(prs, blank, client, insights):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "Strategic Recommendations", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    recs = []
    if isinstance(insights, dict):
        recs = insights.get("recommendations", [])
    if recs:
        y = 1.2
        for i, rec in enumerate(recs[:5], 1):
            if isinstance(rec, dict):
                title = rec.get("title", "")
                desc = rec.get("description", "")
                if len(desc) > 250:
                    desc = desc[:247] + "..."
                text = f"{i}. {title}: {desc}"
            else:
                text = f"{i}. {rec}"
            add_text_box(slide, text, 0.5, y, 12, 0.8,
                         font_size=11, color="1A2B4A")
            y += 0.95
    else:
        add_text_box(slide, "Recommendations will be available once an Anthropic API key is configured in Settings.",
                     0.5, 2, 12, 1, font_size=14, color="888888")
    add_footer(slide, client)


def build_plan_slide(prs, blank, client, insights):
    slide = prs.slides.add_slide(blank)
    add_header_rule(slide)
    add_text_box(slide, "90-Day Plan", 0.5, 0.15, 6, 0.5,
                 font_size=28, bold=True, color="0D1B2E")
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.0), Inches(5.8), Inches(0.5)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = hex_to_rgb("0097B2")
    left_header.line.fill.background()
    add_text_box(slide, "INTERVIEW", 0.7, 1.05, 5, 0.4,
                 font_size=14, bold=True, color="FFFFFF")
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(1.0), Inches(5.8), Inches(0.5)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = hex_to_rgb("0097B2")
    right_header.line.fill.background()
    add_text_box(slide, "BYLINE", 7.2, 1.05, 5, 0.4,
                 font_size=14, bold=True, color="FFFFFF")
    months = []
    if isinstance(insights, dict):
        months = insights.get("months", [])
    y = 1.7
    if not months:
        add_text_box(slide, "90-day plan will be available once an Anthropic API key is configured in Settings.",
                     0.5, 2.0, 12, 1, font_size=14, color="888888")
    for m in months[:3]:
        if not isinstance(m, dict):
            continue
        month_label = m.get("month", "Month")
        add_text_box(slide, month_label, 0.5, y, 5.8, 0.3,
                     font_size=12, bold=True, color="0D1B2E")
        add_text_box(slide, month_label, 7, y, 5.8, 0.3,
                     font_size=12, bold=True, color="0D1B2E")
        y += 0.35
        interviews = m.get("interviews", [])
        bylines = m.get("bylines", [])
        int_y = y
        for topic in interviews[:3]:
            add_text_box(slide, f"• {topic}", 0.7, int_y, 5.5, 0.25,
                         font_size=10, color="1A2B4A")
            int_y += 0.28
        by_y = y
        for topic in bylines[:3]:
            add_text_box(slide, f"• {topic}", 7.2, by_y, 5.5, 0.25,
                         font_size=10, color="1A2B4A")
            by_y += 0.28
        y = max(int_y, by_y) + 0.2
    add_footer(slide, client)


def build_closing_slide(prs, blank, client):
    slide = prs.slides.add_slide(blank)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(6.66), Inches(SLIDE_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb("0097B2")
    shape.line.fill.background()
    add_text_box(slide, "Thank you!", 1, 2.5, 5, 1.5,
                 font_size=44, bold=True, color="FFFFFF")

    # Active DMC logo or text fallback
    admc_logo = load_admc_logo()
    if admc_logo:
        _add_logo_scaled(slide, admc_logo, 7.5, 1.2, max_w=3.0, max_h=1.2)
    else:
        add_text_box(slide, "Active DMC", 7.5, 1.5, 5, 0.5,
                     font_size=24, bold=True, color="0097B2")

    add_text_box(slide, "DIGITAL  •  MARKETING  •  COMMUNICATIONS", 7.5, 2.2, 5, 0.4,
                 font_size=11, color="0097B2")

    # Client logo (bottom-right of the closing slide)
    client_logo = load_logo(client)
    if client_logo:
        _add_logo_scaled(slide, client_logo, 7.5, 5.3, max_w=2.5, max_h=1.0)

    contact_lines = [
        "info@activedmc.com",
        "Dubai, United Arab Emirates",
        "+971 4 XXX XXXX",
        "www.activedmc.com",
    ]
    y = 3.5
    for line in contact_lines:
        add_text_box(slide, line, 7.5, y, 5, 0.3,
                     font_size=12, color="1A2B4A")
        y += 0.4


def build_report(client, period, tracker_data, deliverables, insights, output_path,
                 deliverables_note=None):
    """Build the complete ADMC branded PowerPoint report."""
    tracker_data = tracker_data or {}
    deliverables = deliverables or []
    insights = insights or {}
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    build_cover_slide(prs, blank, client, period)
    build_contents_slide(prs, blank, client, None)

    build_section_divider(prs, blank, "Campaign Overview")
    build_kpi_slide(prs, blank, client, tracker_data, insights.get("campaign_overview", {}))

    build_section_divider(prs, blank, "Coverage Highlights")
    build_highlights_slide(prs, blank, client, tracker_data)

    build_section_divider(prs, blank, "Deliverables Overview")
    build_deliverables_slide(prs, blank, client, deliverables, note=deliverables_note)

    build_section_divider(prs, blank, "Coverage Analysis")
    build_analysis_slide(prs, blank, client, tracker_data, insights.get("coverage_analysis", {}))

    build_section_divider(prs, blank, "Top Publications")
    build_top_publications_slide(prs, blank, client, tracker_data)

    build_section_divider(prs, blank, "Media Relations")
    build_media_relations_slide(prs, blank, client, insights.get("media_relations", {}))

    build_section_divider(prs, blank, "Driving the Narrative")
    build_narrative_slide(prs, blank, client, insights.get("narrative", {}))

    build_section_divider(prs, blank, "Observation / Feedback")
    build_observation_slide(prs, blank, client, insights.get("observation", {}))

    build_section_divider(prs, blank, "Strategic Recommendations")
    build_recommendations_slide(prs, blank, client, insights.get("recommendations", {}))

    build_section_divider(prs, blank, "90-Day Plan")
    build_plan_slide(prs, blank, client, insights.get("ninety_day_plan", {}))

    build_closing_slide(prs, blank, client)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path
